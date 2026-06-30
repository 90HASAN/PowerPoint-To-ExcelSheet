import re
import os
from pptx import Presentation
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side


# ─────────────────────────────────────────────────────────────────────────────
# STEP 1: Extract all text from a slide into a flat list of strings
# ─────────────────────────────────────────────────────────────────────────────

def extract_text_from_slide(slide):
    """Extract all non-empty text from a slide, including tables."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    texts.append(line)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    line = cell.text.strip()
                    if line:
                        texts.append(line)
    return texts


def get_font_size(shape):
    """Get the maximum font size used in a shape (to detect title)."""
    max_size = 0
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    size = run.font.size.pt
                    if size > max_size:
                        max_size = size
    return max_size


def extract_shapes_with_metadata(slide):
    """Extract text with font size info — needed for Type 1 title detection."""
    shapes_data = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        font_size = get_font_size(shape)
        shapes_data.append({
            "text": text,
            "font_size": font_size,
            "left": shape.left,
            "top": shape.top,
        })
    return shapes_data


# ─────────────────────────────────────────────────────────────────────────────
# STEP 2: Detect which type a slide belongs to
# ─────────────────────────────────────────────────────────────────────────────

def detect_slide_type(texts):
    combined = " ".join(texts).upper()

    if re.search(r'DEALER\s+NAME\.', combined) or 'GSTIN' in combined:
        return 'type2'

    if re.search(r'NAME\s*:-', combined) or re.search(r'PH\s*NO\s*:-', combined):
        return 'type3'

    if re.search(r'\bADD\s*:', combined) or re.search(r'\bMOB\s*:', combined):
        return 'type1'

    return 'unknown'


# ─────────────────────────────────────────────────────────────────────────────
# STEP 3: Parsers for each type
# ─────────────────────────────────────────────────────────────────────────────

def clean(s):
    return re.sub(r'\s+', ' ', s).strip() if s else ''


def parse_size(size_str):
    if not size_str:
        return '', '', ''

    matches = list(re.finditer(r'\d+(?:\.\d+)?', size_str))
    if len(matches) < 2:
        return '', '', clean(size_str)

    height, width = matches[0].group(), matches[1].group()

    # Type = letters (plus spaces/hyphens between them) found after the
    # second number, ignoring any leftover unit symbols/punctuation.
    after_second_number = size_str[matches[1].end():]
    letters = re.findall(r'[A-Za-z][A-Za-z\s\-]*', after_second_number)
    size_type = clean(' '.join(letters))

    return height, width, size_type


def parse_type1(slide):
    shapes = extract_shapes_with_metadata(slide)
    texts = [s['text'] for s in shapes]

    name = gst = dl = phone = address = size = ''

    if shapes:
        largest = max(shapes, key=lambda s: s['font_size'])
        if largest['font_size'] > 20:
            name = clean(largest['text'])

    combined = '\n'.join(texts)

    m = re.search(r'ADD\s*[:\-]+\s*(.+)', combined, re.IGNORECASE)
    if m:
        address = clean(m.group(1))

    m = re.search(r'MOB\s*[:\-]+\s*(.+)', combined, re.IGNORECASE)
    if m:
        phone = clean(m.group(1))

    m = re.search(r'SIZE\s*[:\-]+\s*(.+)', combined, re.IGNORECASE)
    if m:
        size = clean(m.group(1))

    m = re.search(r'GST(?:IN)?\s*(?:NO\.?)?\s*[:\-]?\s*([A-Z0-9]{10,})', combined, re.IGNORECASE)
    if m:
        gst = clean(m.group(1))

    m = re.search(r'D\.?L\.?\s*(?:NO\.?)?\s*[:\-]?\s*([A-Z0-9/\-]+)', combined, re.IGNORECASE)
    if m:
        dl = clean(m.group(1))

    return {
        'NAME': name, 'PHONE': phone, 'GST': gst,
        'DL_NO': dl, 'ADDRESS': address, 'SIZE': size
    }


def parse_type2_table(slide):
    name = gst = dl = phone = address = size = ''

    for shape in slide.shapes:
        if not shape.has_table:
            continue
        rows = shape.table.rows
        for row in rows:
            cells = [cell.text.strip() for cell in row.cells]
            combined = ' | '.join(cells)

            m = re.search(r'DEALER\s+NAME\.?\s*(.+?)(?:\s*\||\s{2,}|$)', combined, re.IGNORECASE)
            if m:
                name = clean(m.group(1))

            m = re.search(r'MOB\.?\s*([\d,\s]+)', combined, re.IGNORECASE)
            if m:
                phone = clean(m.group(1))

            m = re.search(r'GSTIN\s*(?:NO\.?)?\s*([A-Z0-9]{10,})', combined, re.IGNORECASE)
            if m:
                gst = clean(m.group(1))

            m = re.search(r'D\.?L\.?\s*NO\.?\s*([A-Z0-9/\-]{3,})', combined, re.IGNORECASE)
            if m:
                val = clean(m.group(1))
                if not re.match(r'^(GSTIN|ADDRESS|SIZE|MOB|DEALER)', val, re.IGNORECASE):
                    dl = val

            m = re.search(r'ADDRESS\.?\s*(.+?)(?:\s*\|)', combined, re.IGNORECASE)
            if m:
                address = clean(m.group(1))

            m = re.search(r'SIZE\.?\s*(.+)', combined, re.IGNORECASE)
            if m:
                size = clean(m.group(1))

    return {
        'NAME': name, 'PHONE': phone, 'GST': gst,
        'DL_NO': dl, 'ADDRESS': address, 'SIZE': size
    }


def parse_type3(slide):
    name = gst = dl = phone = address = size = ''

    shapes_text = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        lines = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
        if not lines:
            continue
        shapes_text.append(lines)

    label_index = None
    for idx, lines in enumerate(shapes_text):
        combined = ' '.join(lines).upper()
        if 'NAME' in combined and ('SIZE' in combined or 'PH NO' in combined or 'ADD' in combined):
            label_index = idx
            break

    if label_index is None:
        combined = '\n'.join([l for lines in shapes_text for l in lines])
        m = re.search(r'NAME\s*[:\-]+\s*(.+)', combined, re.IGNORECASE)
        if m: name = clean(m.group(1))
        m = re.search(r'SIZE\s*[:\-]+\s*(.+)', combined, re.IGNORECASE)
        if m: size = clean(m.group(1))
        m = re.search(r'PH\s*(?:NO\.?)?\s*[:\-]+\s*(.+)', combined, re.IGNORECASE)
        if m: phone = clean(m.group(1))
        m = re.search(r'ADD(?:RESS)?\s*[:\-]+\s*(.+)', combined, re.IGNORECASE)
        if m: address = clean(m.group(1))
        return {'NAME': name, 'PHONE': phone, 'GST': gst, 'DL_NO': dl, 'ADDRESS': address, 'SIZE': size}

    label_lines = shapes_text[label_index]
    field_order = []
    for line in label_lines:
        u = line.upper()
        if re.search(r'NAME', u):       field_order.append('NAME')
        elif re.search(r'SIZE', u):     field_order.append('SIZE')
        elif re.search(r'PH', u):       field_order.append('PHONE')
        elif re.search(r'ADD', u):      field_order.append('ADDRESS')
        else:                           field_order.append(None)

    value_shapes = shapes_text[label_index + 1:]

    result = {'NAME': '', 'SIZE': '', 'PHONE': '', 'ADDRESS': ''}
    for field, val_lines in zip([f for f in field_order if f], value_shapes):
        result[field] = clean(' '.join(val_lines))

    return {
        'NAME':    result.get('NAME', ''),
        'PHONE':   result.get('PHONE', ''),
        'GST':     gst,
        'DL_NO':   dl,
        'ADDRESS': result.get('ADDRESS', ''),
        'SIZE':    result.get('SIZE', '')
    }


# ─────────────────────────────────────────────────────────────────────────────
# STEP 4: Process a single PPT file
# ─────────────────────────────────────────────────────────────────────────────

def process_ppt(ppt_path):
    """Process all slides in a PPT and return list of extracted records."""
    prs = Presentation(ppt_path)
    records = []
    filename = os.path.basename(ppt_path)

    for slide_num, slide in enumerate(prs.slides, 1):
        texts = extract_text_from_slide(slide)
        if not texts:
            continue

        slide_type = detect_slide_type(texts)

        if slide_type == 'type1':
            data = parse_type1(slide)
        elif slide_type == 'type2':
            data = parse_type2_table(slide)
        elif slide_type == 'type3':
            data = parse_type3(slide)
        else:
            results = [parse_type1(slide), parse_type2_table(slide), parse_type3(slide)]
            data = max(results, key=lambda d: sum(1 for v in d.values() if v))

        if data['NAME'] or data['PHONE']:
            data['_file'] = filename
            data['_slide'] = slide_num
            data['_type'] = slide_type
            records.append(data)

    return records


# ─────────────────────────────────────────────────────────────────────────────
# STEP 5: Write to Excel
# ─────────────────────────────────────────────────────────────────────────────

def write_excel(records, output_path):
    """
    output_path may be a filesystem path (str) or a file-like object
    (e.g. io.BytesIO) that openpyxl's Workbook.save() can write to.
    """
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Extracted Data"

    header_fill = PatternFill("solid", fgColor="1F4E79")
    header_font = Font(bold=True, color="FFFFFF", size=11)
    border = Border(
        left=Side(style='thin'), right=Side(style='thin'),
        top=Side(style='thin'), bottom=Side(style='thin')
    )

    headers = ['S.No', 'NAME', 'PHONE', 'GST', 'DL NO', 'ADDRESS', 'HEIGHT', 'WIDTH', 'TYPE']
    col_widths = [6, 35, 20, 20, 20, 45, 12, 12, 20]

    for col, (h, w) in enumerate(zip(headers, col_widths), 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal='center', vertical='center')
        cell.border = border
        ws.column_dimensions[cell.column_letter].width = w

    ws.row_dimensions[1].height = 20

    fill_even = PatternFill("solid", fgColor="EBF3FB")
    fill_odd  = PatternFill("solid", fgColor="FFFFFF")

    for i, rec in enumerate(records, 1):
        row = i + 1
        fill = fill_even if i % 2 == 0 else fill_odd
        height, width, size_type = parse_size(rec.get('SIZE', ''))
        values = [
            i,
            rec.get('NAME', ''),
            rec.get('PHONE', ''),
            rec.get('GST', ''),
            rec.get('DL_NO', ''),
            rec.get('ADDRESS', ''),
            height,
            width,
            size_type,
        ]
        for col, val in enumerate(values, 1):
            cell = ws.cell(row=row, column=col, value=val)
            cell.fill = fill
            cell.border = border
            cell.alignment = Alignment(vertical='center', wrap_text=True)
            if col == 1:
                cell.alignment = Alignment(horizontal='center', vertical='center')

    ws.freeze_panes = 'A2'
    ws.auto_filter.ref = f"A1:I{len(records)+1}"

    wb.save(output_path)
    return wb
